import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_8slide_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    COLOR_BG = RGBColor(15, 23, 42)          # Dark Slate #0F172A
    COLOR_CARD = RGBColor(30, 41, 59)        # Midnight Slate #1E293B
    COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate #334155
    COLOR_PRIMARY = RGBColor(99, 102, 241)   # Indigo Accent #6366F1
    COLOR_CYAN = RGBColor(6, 182, 212)       # Neon Cyan #06B6D4
    COLOR_EMERALD = RGBColor(16, 185, 129)   # Emerald #10B981
    COLOR_ROSE = RGBColor(244, 63, 94)       # Coral Rose #F43F5E
    COLOR_AMBER = RGBColor(245, 158, 11)     # Amber Gold #F59E0B
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# White #F8FAFC
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Muted Gray #94A3B8

    def apply_slide_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_slide_header(slide, tag_text, title_text):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_tag = tf.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_CYAN
        p_tag.font.name = 'Calibri'
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        p_title.font.name = 'Calibri'

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # SLIDE 1: TITLE & STRATEGIC VISION
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide1)
    add_card(slide1, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5), bg_color=COLOR_CARD, border_color=COLOR_PRIMARY)
    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "PRIVACYLENS"
    p0.font.size = Pt(46)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_CYAN

    p1 = tf.add_paragraph()
    p1.text = "Autonomous Digital Footprint & Compliance Control Center for DPDP Act 2023"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN
    p1.space_before = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "A passive Chrome MV3 extension and web control platform empowering users to discover personal data sprawl, understand privacy policies via AI, execute 3-tier consent revocations, and maintain immutable proof of legal compliance."
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.space_before = Pt(15)

    p3 = tf.add_paragraph()
    p3.text = "EXECUTIVE COMMITTEE PRESENTATION • PRELIMINARY BUILD (v1.0)\nStatus: PRODUCTION READY (GREEN) | Law: DPDP Act 2023 Enforcement"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_EMERALD
    p3.space_before = Pt(30)

    # SLIDE 2: THE PROBLEM & COMPLIANCE MANDATE
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide2)
    add_slide_header(slide2, "SLIDE 02 / EXECUTIVE CONTEXT", "The Digital Footprint Crisis & DPDP Act Mandate")

    problems = [
        ("1. Data Sprawl Across 50+ Platforms", "Consumers register across dozens of services leaving email, phone, and name identifiers in legacy corporate databases without ongoing visibility.", COLOR_ROSE),
        ("2. Opaque Consents & Dark Patterns", "Pre-checked checkboxes and 4,000-word terms trick users into granting broad consent for marketing and 3rd-party data monetization.", COLOR_AMBER),
        ("3. Unenforceable DPDP Rights", "DPDP Act 2023 guarantees §6 Revocation and §12 Erasure, but users lack a unified mechanism to exercise rights without visiting 50 sites manually.", COLOR_CYAN),
        ("4. Lack of Audit-Ready Proof", "Without timestamped, verifiable logs of deletion requests, Data Principals cannot prove compliance failure during legal disputes.", COLOR_PRIMARY)
    ]

    for idx, (title, desc, accent) in enumerate(problems):
        row, col = idx // 2, idx % 2
        left, top = Inches(0.8 + col * 5.95), Inches(1.5 + row * 2.7)
        add_card(slide2, left, top, Inches(5.75), Inches(2.5))
        line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.75), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()

        tb = slide2.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25), Inches(5.15), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.space_before = Pt(10)

    # SLIDE 3: PRODUCT VISION & UVP
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide3)
    add_slide_header(slide3, "SLIDE 03 / SOLUTION OVERVIEW", "Product Vision & Unique Value Proposition (UVP)")

    add_card(slide3, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.4))
    tb_l = slide3.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.15), Inches(5.0))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "System Scope & Boundaries"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    items = [
        ("WHAT PRIVACYLENS IS", COLOR_EMERALD),
        ("• Passive event detector via Chrome Extension (MV3)", COLOR_TEXT_MAIN),
        ("• Centralized privacy control center for user digital footprint", COLOR_TEXT_MAIN),
        ("• Realistic 3-Tier Deletion & Consent Revocation Engine", COLOR_TEXT_MAIN),
        ("• Immutable audit trail engine for DPDP compliance proof", COLOR_TEXT_MAIN),
        ("\nWHAT PRIVACYLENS IS NOT", COLOR_ROSE),
        ("• NOT a document vault or DigiLocker storage system", COLOR_TEXT_MUTED),
        ("• NOT for storing Aadhaar, PAN, or personal certificates", COLOR_TEXT_MUTED),
        ("• NOT a password manager or automated form-filler", COLOR_TEXT_MUTED),
        ("• NOT claiming false universal 1-click auto-deletion", COLOR_TEXT_MUTED)
    ]
    for txt, col in items:
        p = tf_l.add_paragraph()
        p.text = txt
        p.font.size = Pt(12)
        p.font.bold = ("WHAT" in txt)
        p.font.color.rgb = col
        p.space_before = Pt(4)

    add_card(slide3, Inches(6.75), Inches(1.5), Inches(5.78), Inches(5.4))
    tb_r = slide3.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.18), Inches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "4 Pillar Value Drivers"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    uvps = [
        ("1. Zero-Effort Auto Discovery", "Extension passively logs DOM form submits & checkboxes live as user browses."),
        ("2. Plain-Language AI Summaries", "Compresses 4,000-word terms into 2 plain-English sentences using Gemini AI."),
        ("3. Honest 3-Tier Action Engine", "Direct API calls, Guided URL drawers, & Legal Notice drafting for 100% reliability."),
        ("4. Audit-Ready DPDP Proof", "Generates timestamped chronological logs proving consent withdrawal & deletion requests.")
    ]
    for title, desc in uvps:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_AMBER
        p_t.space_before = Pt(10)
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    # SLIDE 4: THE FIVE PRODUCT PILLARS
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide4)
    add_slide_header(slide4, "SLIDE 04 / CORE CAPABILITIES", "The Five Core Product Pillars")

    pillars = [
        ("1. DISCOVER", "Passive MV3 Interceptor", "Detects form submits & consent checkboxes automatically without manual entry.", COLOR_CYAN),
        ("2. UNDERSTAND", "AI Policy Intelligence", "Summarizes legal terms via Gemini AI & calculates Privacy Score (0-100).", COLOR_PRIMARY),
        ("3. CONTROL", "3-Tier Action Engine", "Executes Direct API, Guided URL Drawer, or Legal DPDP Notice generator.", COLOR_AMBER),
        ("4. CLEAN UP", "Inactive Account Radar", "Flags stale accounts & high-risk domains for proactive data scrubbing.", COLOR_ROSE),
        ("5. PROVE", "Immutable Audit Engine", "Maintains chronological timeline of consent revocations & request dispatches.", COLOR_EMERALD)
    ]
    card_width = Inches(2.2)
    gap = Inches(0.18)
    for i, (title, sub, body, color) in enumerate(pillars):
        left = Inches(0.8 + i * (2.2 + gap))
        top = Inches(1.6)
        add_card(slide4, left, top, card_width, Inches(5.2))
        header_bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_width, Inches(0.8))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = color
        header_bar.line.fill.background()

        tb_h = slide4.shapes.add_textbox(left, top + Inches(0.1), card_width, Inches(0.6))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.alignment = PP_ALIGN.CENTER
        p_h.font.size = Pt(13)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_TEXT_MAIN

        tb_b = slide4.shapes.add_textbox(left + Inches(0.15), top + Inches(1.0), Inches(1.9), Inches(4.0))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        p_sub = tf_b.paragraphs[0]
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_TEXT_MAIN
        p_body = tf_b.add_paragraph()
        p_body.text = body
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = COLOR_TEXT_MUTED
        p_body.space_before = Pt(12)

    # SLIDE 5: TECHNICAL BREAKTHROUGH: 3-TIER ENGINE
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide5)
    add_slide_header(slide5, "SLIDE 05 / TECHNICAL BREAKTHROUGH", "The Honest 3-Tier Deletion & Revocation Model")

    tiers = [
        ("LEVEL 1 — DIRECT API", "Integrated Partner API Execution", 
         "• Used when target website exposes an official privacy endpoint.\n• Platform calls POST /api/partner/revoke or /delete.\n• Status updates instantly to REVOKED / COMPLETED.\n• Demonstrated via ShopEase partner API backend integration.", COLOR_EMERALD),
        ("LEVEL 2 — GUIDED URL", "Authenticated Redirection Drawer", 
         "• Used when website requires manual login & account navigation.\n• Platform opens guided side-drawer with direct deep link.\n• Provides clear step-by-step UI instructions.\n• Eliminates searching through complex settings pages.", COLOR_AMBER),
        ("LEVEL 3 — LEGAL NOTICE", "DPDP §6 / §12 Legal Generator", 
         "• Used when website lacks public deletion APIs.\n• Generates structured legal demand citing DPDP Act 2023.\n• Triggers pre-filled mailto: dispatch or copyable transcript.\n• Registers request in Request Tracker for statutory SLA monitoring.", COLOR_CYAN)
    ]
    for idx, (title, sub, body, color) in enumerate(tiers):
        left = Inches(0.8 + idx * 3.98)
        top = Inches(1.5)
        add_card(slide5, left, top, Inches(3.78), Inches(5.4))
        bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.78), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tb = slide5.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), Inches(3.28), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_s = tf.add_paragraph()
        p_s.text = sub
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = COLOR_TEXT_MAIN
        p_s.space_before = Pt(4)
        p_b = tf.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = COLOR_TEXT_MUTED
        p_b.space_before = Pt(14)

    # SLIDE 6: SYSTEM ARCHITECTURE & SECURITY
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide6)
    add_slide_header(slide6, "SLIDE 06 / ARCHITECTURE & SECURITY", "End-to-End System Architecture & Zero-Trust Security Mandate")

    add_card(slide6, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.4))
    tb_l = slide6.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.15), Inches(5.0))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "System Integration Flow"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    arch_points = [
        ("1. Chrome MV3 Extension", "Content script listens to DOM form submits & sends sanitized metadata to background worker."),
        ("2. Node.js Express REST Backend", "Receives POST /api/events, executes 3-tier controllers, & communicates with Gemini AI."),
        ("3. Prisma ORM & SQLite Database", "Stores relational records for Users, Websites, Active Consents, Requests, & Audit Logs."),
        ("4. React Web Control Dashboard", "Renders Digital Footprint grid, Privacy Score (0-100), 3-Tier drawers, & Audit timeline.")
    ]
    for t, d in arch_points:
        p_t = tf_l.add_paragraph()
        p_t.text = t
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        p_t.space_before = Pt(8)
        p_d = tf_l.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    add_card(slide6, Inches(6.75), Inches(1.5), Inches(5.78), Inches(5.4), border_color=COLOR_ROSE)
    tb_r = slide6.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.18), Inches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "🛡️ Zero Credential Retention Mandate"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE

    sec_points = [
        ("PASSWORD & PII NON-CAPTURE GUARANTEE", COLOR_AMBER),
        ("• Passwords Ignored: input[type='password'] strictly filtered out.", COLOR_TEXT_MAIN),
        ("• Financial Data Ignored: Credit cards, CVVs, & OTPs never processed.", COLOR_TEXT_MAIN),
        ("• Category Level Metadata: Logs data types ('Email Present') not raw text.", COLOR_TEXT_MAIN),
        ("\nPRIVACY & SECURITY BY DESIGN", COLOR_EMERALD),
        ("• Backend API Key Isolation: Gemini keys stored safely in backend .env.", COLOR_TEXT_MUTED),
        ("• Open Content Script: Execution logic 100% visible in Chrome DevTools.", COLOR_TEXT_MUTED)
    ]
    for txt, col in sec_points:
        p = tf_r.add_paragraph()
        p.text = txt
        p.font.size = Pt(11.5)
        p.font.bold = ("GUARANTEE" in txt or "DESIGN" in txt)
        p.font.color.rgb = col
        p.space_before = Pt(4)

    # SLIDE 7: LIVE DEMO & VERIFICATION HARNESS
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide7)
    add_slide_header(slide7, "SLIDE 07 / LIVE DEMONSTRATION", "Live Verification Harness: ShopEase E-Commerce Demo")

    add_card(slide7, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tb = slide7.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.133), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "100% Controlled End-to-End Live Verification Protocol (10 Steps)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    steps = [
        "1. Launch ShopEase mock e-commerce site (localhost:3000) with Chrome Extension active.",
        "2. User fills out registration form (Name, Email, Phone) & checks Marketing Consent box.",
        "3. Extension DOM listener detects submission; background worker dispatches event payload to REST API.",
        "4. Extension toast triggers: 'Privacy Activity Detected on ShopEase'. Extension popup displays active consents.",
        "5. Launch Central Privacy Dashboard (localhost:5173). Footprint grid renders ShopEase card with Medium risk rating.",
        "6. Digital Privacy Score automatically recalculates (e.g. 85 ➔ 72) due to active marketing consent.",
        "7. Click [AI Policy Summary] to view 2-sentence plain-English legal breakdown powered by Gemini AI.",
        "8. Click [Revoke Marketing Consent] ➔ Tier 1 Direct API executes ➔ Status updates instantly to REVOKED.",
        "9. Click [Request Account Deletion] ➔ Tier 3 Generator drafts DPDP §12 erasure legal notice ➔ Click [Submit Request].",
        "10. Privacy Request Tracker updates status to PENDING RESPONSE & Audit Log timeline registers timestamped proof."
    ]
    for s in steps:
        p_s = tf.add_paragraph()
        p_s.text = s
        p_s.font.size = Pt(11.5)
        p_s.font.color.rgb = COLOR_TEXT_MAIN
        p_s.space_before = Pt(5)

    # SLIDE 8: IMPLEMENTATION ROADMAP & EXECUTION
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide8)
    add_slide_header(slide8, "SLIDE 08 / ROADMAP & EXECUTION", "Implementation Progress, Team Ownership & Roadmap")

    add_card(slide8, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.4))
    tb_l = slide8.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.15), Inches(5.0))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Team Ownership & Status (GREEN)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD

    team_data = [
        ("TM 1 (Team Lead): Integration Architect", "shared/contracts.json, MASTER_DOC, Pitch & Deck", "COMPLETED"),
        ("TM 2 (Extension): Chrome MV3 Eng.", "extension/content.js, background.js, popup/*", "COMPLETED"),
        ("TM 3 (Dashboard): React Frontend", "src/components/*, mockDashboardData.json", "COMPLETED"),
        ("TM 4 (Backend): Express & Prisma", "backend/server.js, schema.prisma, shopease-demo/*", "COMPLETED"),
        ("TM 5 (AI & Intelligence): AI Eng.", "backend/services/aiService.js, privacyScore.js", "COMPLETED")
    ]
    for role, files, stat in team_data:
        p_r = tf_l.add_paragraph()
        p_r.text = role
        p_r.font.size = Pt(12)
        p_r.font.bold = True
        p_r.font.color.rgb = COLOR_TEXT_MAIN
        p_r.space_before = Pt(6)
        p_f = tf_l.add_paragraph()
        p_f.text = f"{files} — [{stat}]"
        p_f.font.size = Pt(10.5)
        p_f.font.color.rgb = COLOR_EMERALD
        p_f.space_before = Pt(1)

    add_card(slide8, Inches(6.75), Inches(1.5), Inches(5.78), Inches(5.4))
    tb_r = slide8.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.18), Inches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "🚀 Future Production Milestones"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    future = [
        ("Phase 2: Automated Email Status Parsing", "Integrate IMAP / Gmail API to automatically parse response emails from data fiduciaries and update Request Tracker status."),
        ("Phase 3: Multi-Language Legal Notices", "Support localized DPDP notice templates across 12 Indian regional languages for broad accessibility."),
        ("Phase 4: Enterprise Data Fiduciary Portal", "Expose lightweight SDK & webhooks for data fiduciaries to adopt Tier 1 Direct API deletion seamlessly.")
    ]
    for title, desc in future:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_AMBER
        p_t.space_before = Pt(10)
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    out_file = "PrivacyLens_8Slide_Executive_Presentation.pptx"
    prs.save(out_file)
    print(f"Successfully generated 8-Slide PowerPoint presentation: {out_file}")

if __name__ == "__main__":
    build_8slide_presentation()
