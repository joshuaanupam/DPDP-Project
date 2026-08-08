import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Color Palette Definitions
    COLOR_BG = RGBColor(15, 23, 42)          # Dark Slate #0F172A
    COLOR_CARD = RGBColor(30, 41, 59)        # Midnight Slate #1E293B
    COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate #334155
    COLOR_PRIMARY = RGBColor(99, 102, 241)   # Indigo Accent #6366F1
    COLOR_CYAN = RGBColor(6, 182, 212)       # Neon Cyan #06B6D4
    COLOR_EMERALD = RGBColor(16, 185, 129)   # Emerald #10B981
    COLOR_ROSE = RGBColor(244, 63, 94)       # Coral Rose #F43F5E
    COLOR_AMBER = RGBColor(245, 158, 11)     # Amber Gold #F59E0B
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# White #F8FAFC
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Muted Gray #94A3B8

    def apply_slide_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_slide_header(slide, tag_text, title_text):
        # Header Container
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Category Tag
        p_tag = tf.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_CYAN
        p_tag.font.name = 'Calibri'
        
        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        p_title.font.name = 'Calibri'

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide1)

    # Hero Card background
    add_card(slide1, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5), bg_color=COLOR_CARD, border_color=COLOR_PRIMARY)

    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "PRIVACYLENS"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_CYAN
    p0.font.name = 'Calibri'

    p1 = tf.add_paragraph()
    p1.text = "Transparent Digital Footprint & Compliance Platform for DPDP Act 2023"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN
    p1.space_before = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "A passive Chrome MV3 extension & privacy control center enabling users to discover personal data sprawl, understand agreements via AI, execute 3-tier deletions, and maintain legal proof of consent decisions."
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.space_before = Pt(15)

    p3 = tf.add_paragraph()
    p3.text = "PROJECT EVALUATION COMMITTEE PRESENTATION — PRELIMINARY BUILD (v1.0)\nTeam PrivacyLens | Date: August 2026 | Status: PRODUCTION READY (GREEN)"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_EMERALD
    p3.space_before = Pt(30)

    # ==========================================
    # SLIDE 2: PROBLEM STATEMENT
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide2)
    add_slide_header(slide2, "01 / Core Motivation", "The Digital Footprint Crisis & The DPDP Act Context")

    # 4 Problem Cards Layout (2x2 Grid)
    problems = [
        ("1. Data Sprawl Across 50+ Web Services", "Users register on dozens of e-commerce, media, and forum sites over time, leaving sensitive identifiers (Name, Email, Phone) scattered across legacy databases without tracking.", COLOR_ROSE),
        ("2. Invisible & Opaque Consents", "Pre-checked boxes and 40-page legal terms trick users into granting broad consent for marketing, tracking, and 3rd-party data sharing with zero visibility.", COLOR_AMBER),
        ("3. Unenforceable Deletion Rights", "The DPDP Act 2023 guarantees §6 Revocation and §12 Erasure, but users lack simple, consolidated tools to enforce these rights across fragmented web platforms.", COLOR_CYAN),
        ("4. Absence of Audit-Ready Evidence", "When users request account deletion via email or web forms, organizations rarely provide verifiable logs, leaving individuals without proof of legal compliance.", COLOR_PRIMARY)
    ]

    for idx, (title, desc, accent) in enumerate(problems):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.5 + row * 2.7)
        card = add_card(slide2, left, top, Inches(5.75), Inches(2.5))
        
        # Accent indicator line on top of card
        line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.75), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()

        tb = slide2.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25), Inches(5.15), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.space_before = Pt(10)

    # ==========================================
    # SLIDE 3: PRODUCT VISION & UVP
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide3)
    add_slide_header(slide3, "02 / Executive Summary", "Product Vision & Unique Value Proposition (UVP)")

    # Left Box: What PrivacyLens IS vs IS NOT
    add_card(slide3, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.4))
    tb_left = slide3.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.15), Inches(5.0))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "System Boundary & Scope"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    is_items = [
        ("WHAT PRIVACYLENS IS", COLOR_EMERALD),
        ("• Centralized Privacy Control Layer for personal footprint", COLOR_TEXT_MAIN),
        ("• Passive event detector via Chrome Extension (MV3)", COLOR_TEXT_MAIN),
        ("• Realistic 3-Tier Deletion & Revocation Engine", COLOR_TEXT_MAIN),
        ("• Audit & evidence engine for DPDP compliance proof", COLOR_TEXT_MAIN),
        ("\nWHAT PRIVACYLENS IS NOT", COLOR_ROSE),
        ("• NOT a document vault or DigiLocker replacement", COLOR_TEXT_MUTED),
        ("• NOT for storing Aadhaar, PAN, or personal certificates", COLOR_TEXT_MUTED),
        ("• NOT a password manager or automated form-filler", COLOR_TEXT_MUTED),
        ("• NOT claiming false 1-click universal auto-deletion", COLOR_TEXT_MUTED),
    ]

    for text, col in is_items:
        p = tf_l.add_paragraph()
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = ("WHAT" in text)
        p.font.color.rgb = col
        p.space_before = Pt(4)

    # Right Box: 4 Pillars of UVP
    add_card(slide3, Inches(6.75), Inches(1.5), Inches(5.78), Inches(5.4))
    tb_right = slide3.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.18), Inches(5.0))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Four Pillars of UVP"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    uvps = [
        ("1. Zero-Effort Automatic Discovery", "Passive background extension logs form submit metadata live as users navigate the web."),
        ("2. Plain-Language AI Summarization", "Converts multi-page terms of service into 2-sentence actionable summaries via Gemini AI."),
        ("3. Honest 3-Tier Action Engine", "Combines Direct API calls, Guided URLs, and Formatted Legal Request drafting for real-world reliability."),
        ("4. Audit-Ready DPDP Proof", "Generates immutable chronological logs with timestamps proving consent revocation and erasure submission.")
    ]

    for title, desc in uvps:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_AMBER
        p_t.space_before = Pt(10)

        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    # ==========================================
    # SLIDE 4: THE FIVE PRODUCT PILLARS
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide4)
    add_slide_header(slide4, "03 / Product Architecture", "The Five Core Product Pillars")

    pillars = [
        ("1. DISCOVER", "Passive Chrome MV3 Interceptor", "Detects form submits & consent checkboxes without manual data entry.", COLOR_CYAN),
        ("2. UNDERSTAND", "AI Policy & Risk Intelligence", "Summarizes legal terms via Gemini API & calculates Privacy Score (0-100).", COLOR_PRIMARY),
        ("3. CONTROL", "3-Tier Action Engine", "Executes Direct API, Guided Drawer, or Legal Notice generation.", COLOR_AMBER),
        ("4. CLEAN UP", "Inactive Account Radar", "Flags stale accounts & high-risk domains for proactive data scrubbing.", COLOR_ROSE),
        ("5. PROVE", "Immutable Audit Engine", "Maintains chronological audit trail of all consent revocations & requests.", COLOR_EMERALD)
    ]

    card_width = Inches(2.2)
    gap = Inches(0.18)

    for i, (title, sub, body, color) in enumerate(pillars):
        left = Inches(0.8 + i * (2.2 + gap))
        top = Inches(1.6)
        card = add_card(slide4, left, top, card_width, Inches(5.2))

        # Top Accent Header
        header_bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_width, Inches(0.8))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = color
        header_bar.line.fill.background()

        tb_h = slide4.shapes.add_textbox(left, top + Inches(0.1), card_width, Inches(0.6))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.alignment = PP_ALIGN.CENTER
        p_h.font.size = Pt(13)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_TEXT_MAIN

        tb_b = slide4.shapes.add_textbox(left + Inches(0.15), top + Inches(1.0), Inches(1.9), Inches(4.0))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True

        p_sub = tf_b.paragraphs[0]
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_TEXT_MAIN

        p_body = tf_b.add_paragraph()
        p_body.text = body
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = COLOR_TEXT_MUTED
        p_body.space_before = Pt(12)

    # ==========================================
    # SLIDE 5: 3-TIER DELETION ENGINE
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide5)
    add_slide_header(slide5, "04 / Technical Breakthrough", "The Honest 3-Tier Deletion & Revocation Model")

    tiers = [
        ("LEVEL 1 — DIRECT AUTOMATION", "Partner API Endpoint Integration", 
         "• Used when target website exposes an official privacy API.\n• Platform executes POST /api/partner/revoke or /delete.\n• Status updates instantaneously to REVOKED / COMPLETED.\n• Fully validated in demo via ShopEase backend integration.", COLOR_EMERALD),
        ("LEVEL 2 — GUIDED ACTION", "Authenticated Redirection Drawer", 
         "• Used when site requires manual login & account navigation.\n• Platform opens a guided side-drawer with direct deep link.\n• Provides step-by-step UI directions for deletion.\n• Eliminates searching through complex settings menus.", COLOR_AMBER),
        ("LEVEL 3 — LEGAL REQUEST DRAFTING", "DPDP §6 / §12 Notice Generator", 
         "• Used when target website lacks public deletion APIs.\n• Generates structured legal demand citing DPDP Act 2023.\n• Triggers pre-filled mailto: or copyable legal transcript.\n• Registers request in Privacy Request Tracker for SLA tracking.", COLOR_CYAN)
    ]

    for idx, (title, sub, body, color) in enumerate(tiers):
        left = Inches(0.8 + idx * 3.98)
        top = Inches(1.5)
        card = add_card(slide5, left, top, Inches(3.78), Inches(5.4))

        # Accent Bar
        bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.78), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tb = slide5.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), Inches(3.28), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color

        p_s = tf.add_paragraph()
        p_s.text = sub
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = COLOR_TEXT_MAIN
        p_s.space_before = Pt(4)

        p_b = tf.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = COLOR_TEXT_MUTED
        p_b.space_before = Pt(14)

    # ==========================================
    # SLIDE 6: SYSTEM ARCHITECTURE
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide6)
    add_slide_header(slide6, "05 / System Architecture", "End-to-End System Integration Flow")

    # Diagram Boxes (4 Horizontal Steps with Flow Arrows)
    components = [
        ("BROWSER EXTENSION\n(Manifest V3)", "• Content Script DOM Listener\n• Background Service Worker\n• Zero Password Security Mask\n• Extension Popup View", COLOR_CYAN),
        ("NODE.JS / EXPRESS\nREST BACKEND", "• Auth & Security Middleware\n• Event Ingestion Controller\n• 3-Tier Action Engine\n• Gemini AI Policy Parser", COLOR_PRIMARY),
        ("PRISMA ORM &\nSQLITE DATABASE", "• User Footprint Table\n• Active Consents Table\n• Privacy Requests Table\n• Immutable Audit Log", COLOR_AMBER),
        ("REACT CONTROL\nDASHBOARD", "• Digital Footprint Grid\n• Dynamic Privacy Score (0-100)\n• 3-Tier Execution Drawers\n• Audit Log Timeline", COLOR_EMERALD)
    ]

    for idx, (title, details, color) in enumerate(components):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)
        card = add_card(slide6, left, top, Inches(2.6), Inches(4.5))

        # Title Card Header
        t_bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.6), Inches(0.9))
        t_bar.fill.solid()
        t_bar.fill.fore_color.rgb = color
        t_bar.line.fill.background()

        tb_t = slide6.shapes.add_textbox(left, top + Inches(0.1), Inches(2.6), Inches(0.7))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        p = tf_t.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

        tb_d = slide6.shapes.add_textbox(left + Inches(0.15), top + Inches(1.1), Inches(2.3), Inches(3.2))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = details
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

        # Add Connector Arrow between components
        if idx < 3:
            arr_left = left + Inches(2.6)
            arr = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arr_left + Inches(0.05), top + Inches(2.0), Inches(0.25), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_TEXT_MUTED
            arr.line.fill.background()

    # ==========================================
    # SLIDE 7: BROWSER EXTENSION MV3
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide7)
    add_slide_header(slide7, "06 / Extension Engineering", "Browser Extension Architecture & Security Mandate")

    # Left Side: Extension Technical Workflow
    add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.4))
    tb_l = slide7.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.15), Inches(5.0))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "Extension Execution Pipeline"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    ext_steps = [
        ("1. Passive DOM Interception (content.js)", "Attaches non-blocking listeners to HTML form submit events across all visited web domains."),
        ("2. Form & Consent Inspection", "Identifies input fields (Email, Phone, Name) and checkboxes containing terms/marketing keywords."),
        ("3. Payload Dispatch (background.js)", "Sends sanitized metadata payload via async HTTP POST /api/events to central backend."),
        ("4. Extension Popup View (popup.jsx)", "Displays instant domain privacy status, detected fields, and active consents to the user.")
    ]

    for title, desc in ext_steps:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        p_t.space_before = Pt(8)

        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    # Right Side: Strict Security & Non-Capture Mandate
    add_card(slide7, Inches(6.75), Inches(1.5), Inches(5.78), Inches(5.4), border_color=COLOR_ROSE)
    tb_r = slide7.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.18), Inches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "🛡️ Zero Credential Retention Mandate"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE

    sec_rules = [
        ("STRICT EXCLUSION RULES", COLOR_AMBER),
        ("• Passwords Completely Ignored: input[type='password'] strictly excluded.", COLOR_TEXT_MAIN),
        ("• Payment Cards Ignored: Credit card, CVV, and banking fields stripped out.", COLOR_TEXT_MAIN),
        ("• OTP & Auth Tokens Ignored: Security verification codes never processed.", COLOR_TEXT_MAIN),
        ("\nPRIVACY BY DESIGN", COLOR_CYAN),
        ("• Category Level Metadata: Logs data types ('Email Present') rather than raw text.", COLOR_TEXT_MUTED),
        ("• Environment Isolation: API keys isolated exclusively in backend .env.", COLOR_TEXT_MUTED),
        ("• Transparent Execution: Open content script execution visible in DevTools.", COLOR_TEXT_MUTED)
    ]

    for text, color in sec_rules:
        p = tf_r.add_paragraph()
        p.text = text
        p.font.size = Pt(11.5)
        p.font.bold = ("STRICT" in text or "PRIVACY" in text)
        p.font.color.rgb = color
        p.space_before = Pt(4)

    # ==========================================
    # SLIDE 8: AI POLICY INTELLIGENCE & PRIVACY SCORE
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide8)
    add_slide_header(slide8, "07 / Privacy Intelligence", "AI Policy Summarizer & Privacy Score Algorithm")

    # Left Card: AI Policy Summarizer
    add_card(slide8, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.4))
    tb_l = slide8.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.15), Inches(5.0))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "🤖 Gemini API Policy Summarizer"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    ai_bullets = [
        ("Problem Solved", "Privacy policies average 4,000 words of legal jargon that no consumer reads."),
        ("Gemini AI Prompt Engineering", "Extracts key data practices into 2 plain-English sentences: (1) What data is collected, (2) How data is shared/sold."),
        ("Rule-Based Graceful Fallback", "If Gemini API key is missing or offline, system switches to structured rule-based keyword extraction without breaking UI.")
    ]

    for title, desc in ai_bullets:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_AMBER
        p_t.space_before = Pt(10)

        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    # Right Card: Privacy Score Formula
    add_card(slide8, Inches(6.75), Inches(1.5), Inches(5.78), Inches(5.4))
    tb_r = slide8.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.18), Inches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "📊 Dynamic Privacy Score Equation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD

    p_eq = tf_r.add_paragraph()
    p_eq.text = "Score = 100 - (ActiveMarketing × 5) - (HighRiskSites × 10) + (RevokedConsents × 3)"
    p_eq.font.size = Pt(12)
    p_eq.font.bold = True
    p_eq.font.color.rgb = COLOR_CYAN
    p_eq.space_before = Pt(12)

    score_rules = [
        ("Base Score: 100", "Starting score for clean digital footprint."),
        ("Marketing Deductions (-5 pts)", "Subtracted for each active promotional / 3rd party consent."),
        ("High Risk Site Penalty (-10 pts)", "Deducted for accounts on unverified or data-broker domains."),
        ("Revocation Bonus (+3 pts)", "Rewarded immediately when user revokes consent or submits erasure.")
    ]

    for title, desc in score_rules:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_MAIN
        p_t.space_before = Pt(8)

        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    # ==========================================
    # SLIDE 9: CENTRAL DASHBOARD & AUDIT PROOF
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide9)
    add_slide_header(slide9, "08 / User Interface & Proof", "Central Privacy Dashboard & Audit Evidence Engine")

    dash_features = [
        ("Digital Footprint Grid", "Card-based interface summarizing all connected websites, active data categories (Email, Phone, Name), risk tags, and deletion tiers.", COLOR_CYAN),
        ("Website Detail Modal", "Granular breakdown of specific domain agreements, displaying plain-English AI summaries and 3-tier execution triggers.", COLOR_PRIMARY),
        ("Privacy Request Tracker", "Monitors lifecycle of deletion requests (CREATED ➔ SUBMITTED ➔ AWAITING RESPONSE ➔ COMPLETED) with legal deadline indicators.", COLOR_AMBER),
        ("Immutable Audit Log Timeline", "Chronological evidence timeline recording exact ISO timestamps for event discovery, consent revocations, and DPDP request submissions.", COLOR_EMERALD)
    ]

    for idx, (title, desc, color) in enumerate(dash_features):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.5 + row * 2.7)
        add_card(slide9, left, top, Inches(5.75), Inches(2.5))

        bar = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.75), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tb = slide9.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25), Inches(5.15), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.space_before = Pt(10)

    # ==========================================
    # SLIDE 10: DEMO HARNESS ("ShopEase")
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide10)
    add_slide_header(slide10, "09 / Verification & Demo", "Live Verification Harness: ShopEase Demo Site")

    add_card(slide10, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tb = slide10.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.133), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "100% Reliable Live Demonstration Protocol (14-Step Flow)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    steps = [
        "1. Open ShopEase e-commerce demo store (localhost:3000) with Chrome Extension active.",
        "2. Complete registration form (Name, Email, Phone) and check Marketing Consent box.",
        "3. Extension content script passively captures submit metadata; background worker posts event to REST API.",
        "4. Extension toast displays: 'Privacy Activity Detected on ShopEase'. Extension popup reflects active consents.",
        "5. Launch Web Dashboard (localhost:5173). Footprint grid renders ShopEase card with Medium risk rating.",
        "6. Observe Digital Privacy Score drop from 85 ➔ 72 due to active marketing consent.",
        "7. Trigger 'AI Policy Summary' modal to view plain-English policy breakdown.",
        "8. Click [Revoke Marketing Consent] ➔ Tier 1 Direct API executes ➔ Consent status updates instantly to REVOKED.",
        "9. Click [Request Account Deletion] ➔ Tier 3 Legal Generator opens ➔ Click [Submit Legal Request].",
        "10. Verify Privacy Request Tracker updates status to PENDING RESPONSE & Audit Log records exact timestamped proof."
    ]

    for s in steps:
        p_s = tf.add_paragraph()
        p_s.text = s
        p_s.font.size = Pt(11.5)
        p_s.font.color.rgb = COLOR_TEXT_MAIN
        p_s.space_before = Pt(5)

    # ==========================================
    # SLIDE 11: TEAM ROLES & STATUS
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide11)
    add_slide_header(slide11, "10 / Team Allocation", "Implementation Progress & File Ownership Matrix")

    # Table of Team Members
    rows = 6
    cols = 4
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(11.733)
    height = Inches(5.4)

    table_shape = slide11.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(4.633)
    table.columns[3].width = Inches(2.5)

    headers = ["Role", "Owner", "Key Files & Responsibility", "Status"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN

    matrix_data = [
        ("TM 1: Team Lead", "Integration Architect", "shared/contracts.json, MASTER_DOC, Pitch & Integration", "COMPLETED"),
        ("TM 2: Extension Eng.", "Chrome MV3 Developer", "extension/content.js, background.js, popup/*", "COMPLETED"),
        ("TM 3: Dashboard Eng.", "React Frontend Developer", "src/components/*, mockDashboardData.json", "COMPLETED"),
        ("TM 4: Backend Eng.", "Express & Prisma Developer", "backend/server.js, schema.prisma, shopease-demo/*", "COMPLETED"),
        ("TM 5: AI & Privacy Eng.", "Intelligence Engineer", "backend/services/aiService.js, privacyScore.js", "COMPLETED")
    ]

    for row_idx, data in enumerate(matrix_data, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10.5)
            if col_idx == 3:
                p.font.bold = True
                p.font.color.rgb = COLOR_EMERALD
            else:
                p.font.color.rgb = COLOR_TEXT_MUTED if col_idx == 2 else COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 12: ROADMAP & Q&A DEFENSE
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide12)
    add_slide_header(slide12, "11 / Future Roadmap", "Strategic Roadmap & Committee Q&A Defense")

    # Left: Roadmap Milestones
    add_card(slide12, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.4))
    tb_l = slide12.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.15), Inches(5.0))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "🚀 Future Production Milestones"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    roadmap = [
        ("Phase 2: Automated Email Parsing", "Integrate IMAP / Gmail API to automatically parse response emails from data fiduciaries and update status in Request Tracker."),
        ("Phase 3: Multi-Language Notices", "Support localized DPDP notice templates across 12 Indian regional languages for broad accessibility."),
        ("Phase 4: Enterprise Partner Portal", "Expose webhooks & SDK for data fiduciaries to implement Tier 1 Direct API deletion seamlessly.")
    ]

    for title, desc in roadmap:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_AMBER
        p_t.space_before = Pt(10)

        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(2)

    # Right: Committee Q&A Defense Box
    add_card(slide12, Inches(6.75), Inches(1.5), Inches(5.78), Inches(5.4), border_color=COLOR_PRIMARY)
    tb_r = slide12.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.18), Inches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "❓ Committee Q&A Readiness"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    qas = [
        ("Q: How do you handle password safety?", "A: Password inputs (type='password') and credit card fields are strictly filtered at content script layer and never leave the DOM."),
        ("Q: Why not auto-delete on all sites?", "A: Universal auto-deletion is impossible due to auth & CSRF. Our 3-tier engine provides an honest, reliable alternative."),
        ("Q: How does this align with DPDP 2023?", "A: Directly enforces Data Principal rights under §6 (Consent Withdrawal) & §12 (Data Erasure) with timestamped audit logs.")
    ]

    for q, a in qas:
        p_q = tf_r.add_paragraph()
        p_q.text = q
        p_q.font.size = Pt(11.5)
        p_q.font.bold = True
        p_q.font.color.rgb = COLOR_EMERALD
        p_q.space_before = Pt(10)

        p_a = tf_r.add_paragraph()
        p_a.text = a
        p_a.font.size = Pt(10.5)
        p_a.font.color.rgb = COLOR_TEXT_MUTED
        p_a.space_before = Pt(2)

    # Output path
    out_file = "PrivacyLens_Preliminary_Presentation.pptx"
    prs.save(out_file)
    print(f"Successfully generated PowerPoint presentation: {out_file}")

if __name__ == "__main__":
    build_presentation()
